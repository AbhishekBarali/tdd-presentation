#!/usr/bin/env python3
"""Build tdd-presentation.pptx from the design of index.html.

The web deck is a fixed 1280x720 canvas. This script rebuilds the same ten
slides as a native PowerPoint file: real text boxes and real shapes, so the
deck stays editable in PowerPoint / Keynote / LibreOffice. Nothing is a
screenshot.

Run:  python tools/build_pptx.py
Out:  tdd-presentation.pptx
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------- design tokens

PAPER = RGBColor(0xF7, 0xF5, 0xF0)
INK = RGBColor(0x19, 0x19, 0x17)
INK_SOFT = RGBColor(0x5B, 0x57, 0x50)
INK_FAINT = RGBColor(0x8B, 0x86, 0x7C)
ACCENT = RGBColor(0xA8, 0x45, 0x2C)
RULE = RGBColor(0xDC, 0xD9, 0xD2)  # flattened rgba(25,25,23,.14) on paper
MARK_FAINT = RGBColor(0xD5, 0xD2, 0xCB)  # the .16-opacity title ornament

DISPLAY = "Georgia"      # stands in for Instrument Serif; ships everywhere
SANS = "Segoe UI"        # stands in for Inter
MONO = "Consolas"        # stands in for JetBrains Mono

STAGE_W, STAGE_H = 1280, 720
PAD_L, PAD_R, PAD_T = 92, 92, 76
CONTENT_W = STAGE_W - PAD_L - PAD_R

TITLE = "Test-Driven Development"
TOTAL = 10


def px(v):
    """Canvas pixels -> EMU, on the 1280x720 = 13.333x7.5in stage."""
    return Emu(int(round(Inches(v / 96.0))))


def flat(shape):
    """Strip the theme shape style so rules and dots render with no shadow.

    An empty <a:effectLst/> alone is not enough: the default autoshape carries
    a <p:style> with an effectRef, and LibreOffice still draws the shadow from
    it. Removing the style element is what actually gives flat ink.
    """
    shape.shadow.inherit = False
    style = shape._element.find(
        "{http://schemas.openxmlformats.org/presentationml/2006/main}style")
    if style is not None:
        shape._element.remove(style)
    return shape


# ------------------------------------------------------------------- primitives


def textbox(slide, x, y, w, h):
    box = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    return tf


def para(tf, first=False, space_before=0, space_after=0, line=None,
         align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if line is not None:
        p.line_spacing = line
    p.alignment = align
    return p


def run(p, text, size, color=INK_SOFT, font=SANS, bold=False, italic=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.name = font
    r.font.bold = bold
    r.font.italic = italic
    return r


def simple(slide, x, y, w, h, text, size, color=INK_SOFT, font=SANS,
           line=1.5, bold=False, align=PP_ALIGN.LEFT):
    """One block of text, newlines become paragraphs."""
    tf = textbox(slide, x, y, w, h)
    for i, chunk in enumerate(text.split("\n")):
        p = para(tf, first=(i == 0), line=line, align=align)
        run(p, chunk, size, color, font, bold)
    return tf


def rule(slide, x, y, w, thickness=1, color=RULE):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(x), px(y), px(w),
                                 px(thickness))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return flat(bar)


def vrule(slide, x, y, h, color=RULE):
    return rule(slide, x, y, 1, thickness=h, color=color)


def line(slide, x1, y1, x2, y2, color=INK_FAINT, width=1.6):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, px(x1), px(y1),
                                   px(x2), px(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width * 0.75)
    return flat(c)


def ring(slide, cx, cy, r, color=INK_FAINT, width=1.6):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, px(cx - r), px(cy - r),
                               px(2 * r), px(2 * r))
    c.fill.background()
    c.line.color.rgb = color
    c.line.width = Pt(width * 0.75)
    return flat(c)


def dot(slide, cx, cy, r, color=INK):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, px(cx - r), px(cy - r),
                               px(2 * r), px(2 * r))
    c.fill.solid()
    c.fill.fore_color.rgb = color
    c.line.fill.background()
    return flat(c)


def label(slide, cx, cy, text, size, color=INK, font=DISPLAY,
          align=PP_ALIGN.CENTER, w=260, h=40):
    tf = textbox(slide, cx - w / 2.0, cy - h / 2.0, w, h)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = para(tf, first=True, line=1.1, align=align)
    run(p, text, size, color, font)
    return tf


def dash(slide, x, y, w=14, color=ACCENT):
    return rule(slide, x, y, w, thickness=2, color=color)


# ------------------------------------------------------------------ slide chrome


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = PAPER
    return slide


def chrome(slide, index):
    """Running head, slide counter and progress bar, as in the web deck."""
    if index > 0:
        simple(slide, PAD_L, STAGE_H - 44 - 18, 400, 24, TITLE, 11, INK_FAINT,
               line=1.0)
        simple(slide, STAGE_W - PAD_R - 200, STAGE_H - 44 - 18, 200, 24,
               "%d / %d" % (index + 1, TOTAL), 11, INK_FAINT, line=1.0,
               align=PP_ALIGN.RIGHT)
    rule(slide, 0, STAGE_H - 4, (index + 1) / float(TOTAL) * STAGE_W,
         thickness=4, color=ACCENT)


def heading(slide, text, y=PAD_T, size=36):
    simple(slide, PAD_L, y, CONTENT_W, 90, text, size, INK, DISPLAY, line=1.05)


# -------------------------------------------------------------------- ornaments


def loop_figure(slide, cx, cy, r, labels=True):
    """The red -> green -> refactor loop from the .fig-loop SVG."""
    ring(slide, cx, cy, r, INK_FAINT, 1.6)
    dot(slide, cx, cy - r, 9, ACCENT)                       # 12 o'clock
    dot(slide, cx + r * 0.866, cy + r * 0.5, 9, INK)        #  4 o'clock
    dot(slide, cx - r * 0.866, cy + r * 0.5, 9, INK)        #  8 o'clock
    if labels:
        label(slide, cx, cy - r - 30, "Red", 20, ACCENT)
        label(slide, cx + r * 0.866 + 74, cy + r * 0.5 + 8, "Green", 20, INK)
        label(slide, cx - r * 0.866 - 78, cy + r * 0.5 + 8, "Refactor", 20, INK)
        label(slide, cx, cy - 12, "one small step,", 12, INK_FAINT, SANS)
        label(slide, cx, cy + 12, "then around again", 12, INK_FAINT, SANS)


def cross_badge(slide, cx, cy, r=20):
    ring(slide, cx, cy, r, ACCENT, 1.6)
    line(slide, cx - 8, cy - 8, cx + 8, cy + 8, ACCENT, 2.4)
    line(slide, cx + 8, cy - 8, cx - 8, cy + 8, ACCENT, 2.4)


def check_badge(slide, cx, cy, r=20):
    ring(slide, cx, cy, r, INK, 1.6)
    line(slide, cx - 9, cy, cx - 2, cy + 7, INK, 2.4)
    line(slide, cx - 2, cy + 7, cx + 11, cy - 8, INK, 2.4)


def run_figure(slide, x, y):
    """fail -> pass strip from the .fig-run SVG."""
    cross_badge(slide, x + 27, y + 34)
    simple(slide, x + 62, y + 16, 220, 26, "The test fails", 15, INK, line=1.1)
    simple(slide, x + 62, y + 42, 240, 24, "the code does not exist yet", 12,
           INK_FAINT, line=1.1)
    line(slide, x + 282, y + 34, x + 316, y + 34, INK_FAINT, 1.6)
    line(slide, x + 310, y + 29, x + 318, y + 34, INK_FAINT, 1.6)
    line(slide, x + 310, y + 39, x + 318, y + 34, INK_FAINT, 1.6)
    check_badge(slide, x + 360, y + 34)
    simple(slide, x + 395, y + 16, 220, 26, "The test passes", 15, INK, line=1.1)
    simple(slide, x + 395, y + 42, 240, 24, "two lines of code later", 12,
           INK_FAINT, line=1.1)


def stack_figure(slide, x, y):
    """The growing safety net from the .fig-stack SVG: one new tick per step."""
    for step in range(1, 6):
        row_y = y + (step - 1) * 46
        simple(slide, x, row_y, 70, 24, "step %d" % step, 11, INK_FAINT,
               line=1.0)
        for k in range(step):
            colour = ACCENT if k == step - 1 else INK_FAINT
            tick = textbox(slide, x + 72 + k * 36, row_y - 4, 34, 30)
            p = para(tick, first=True, line=1.0, align=PP_ALIGN.CENTER)
            run(p, "\u2713", 15, colour, SANS, bold=True)


# ----------------------------------------------------------------------- slides


def slide_1(prs):
    s = new_slide(prs)
    # the faint ornament that sits behind the title
    ring(s, 832 + 172, 196 + 172, 118, MARK_FAINT, 2)
    dot(s, 832 + 172, 196 + 54, 9, MARK_FAINT)
    dot(s, 832 + 274, 196 + 231, 9, MARK_FAINT)
    dot(s, 832 + 70, 196 + 231, 9, MARK_FAINT)

    simple(s, PAD_L, 232, 700, 210, "Test-Driven\nDevelopment", 64, INK,
           DISPLAY, line=0.98)
    simple(s, PAD_L, 452, 700, 40, "Unit 8.2 \u00b7 Software Testing", 20,
           INK_SOFT, line=1.2)

    rule(s, PAD_L, 616, CONTENT_W)
    simple(s, PAD_L, 636, 560, 26,
           "Software Engineering \u00b7 CSC364 \u00b7 Semester VI", 12,
           INK_FAINT, line=1.2)
    simple(s, STAGE_W - PAD_R - 560, 636, 560, 26,
           "Sommerville, Software Engineering, 10th ed.", 12, INK_FAINT,
           line=1.2, align=PP_ALIGN.RIGHT)
    chrome(s, 0)


def slide_2(prs):
    s = new_slide(prs)
    simple(s, PAD_L, 190, 900, 260,
           "Write the test first.\nThen write the code\nthat passes it.", 40,
           INK, DISPLAY, line=1.08)
    simple(s, PAD_L, 452, 640, 140,
           "That is the whole idea. Testing is not a step at the end. "
           "The test and the code grow together, a few lines at a time.", 17,
           INK_SOFT, line=1.45)
    chrome(s, 1)


def slide_3(prs):
    s = new_slide(prs)
    heading(s, "The cycle")

    steps = [
        ("Pick one small piece of work",
         "Something you can write in a few lines."),
        ("Write a test for it",
         "An automatic test, run by a tool, not by hand."),
        ("Run the test and watch it fail",
         "The code is not written yet, so it must fail."),
        ("Write just enough code to pass",
         "Nothing extra. Only what the test asks for."),
        ("Run every test again, then take the next piece",
         "Old tests must still pass before you move on."),
    ]
    y = 186
    for n, (head, sub) in enumerate(steps, start=1):
        simple(s, PAD_L, y - 4, 44, 44, str(n), 22, ACCENT, DISPLAY, line=1.0)
        simple(s, PAD_L + 52, y, 560, 34, head, 17, INK, line=1.25)
        simple(s, PAD_L + 52, y + 28, 560, 30, sub, 13, INK_SOFT, line=1.25)
        y += 92

    loop_figure(s, 990, 400, 118)
    chrome(s, 2)


def slide_4(prs):
    s = new_slide(prs)
    heading(s, "Three words to remember it by")

    phases = [
        ("Red", "The new test fails. Good. It proves the test really checks "
                "something.", ACCENT),
        ("Green", "You write the simplest code that makes the test pass.", INK),
        ("Refactor", "You tidy the code up. The tests tell you if you broke "
                     "anything.", INK),
    ]
    col_w = CONTENT_W / 3.0
    for i, (name, body, colour) in enumerate(phases):
        x = PAD_L + i * col_w
        rule(s, x, 198, col_w - 34, thickness=2, color=INK)
        if i:
            vrule(s, x - 17, 198, 250)
        simple(s, x, 228, col_w - 40, 60, name, 30, colour, DISPLAY, line=1.05)
        simple(s, x, 300, col_w - 40, 150, body, 15, INK_SOFT, line=1.45)

    simple(s, PAD_L, 486, 820, 60,
           "Then you go around again. Each loop is minutes long, not days.", 16,
           INK_SOFT, line=1.45)
    chrome(s, 3)


def slide_5(prs):
    s = new_slide(prs)
    heading(s, "A small example")

    col_w = (CONTENT_W - 44) / 2.0
    code = [
        (PAD_L, "Step 1 \u2014 the test, written first",
         [[("# a 10% discount on any price", INK_FAINT)],
          [("def ", ACCENT), ("test_discount():", INK)],
          [("    assert ", ACCENT), ("discount(1000) == 100", INK)],
          [("    assert ", ACCENT), ("discount(0) == 0", INK)]]),
        (PAD_L + col_w + 44, "Step 2 \u2014 the code, written after",
         [[("def ", ACCENT), ("discount(price):", INK)],
          [("    return ", ACCENT), ("price * 0.10", INK)]]),
    ]
    for x, head, lines in code:
        simple(s, x, 184, col_w, 26, head, 13, INK, bold=True, line=1.2)
        rule(s, x, 214, col_w)
        tf = textbox(s, x, 234, col_w, 160)
        for i, parts in enumerate(lines):
            p = para(tf, first=(i == 0), line=1.5)
            for text, colour in parts:
                run(p, text, 14, colour, MONO)

    simple(s, PAD_L, 404, 760, 90,
           "Run the test before writing discount() and it fails. Write those "
           "two lines and it passes. That pass is now permanent: it runs again "
           "every time anyone touches this code.", 16, INK_SOFT, line=1.45)
    run_figure(s, PAD_L, 528)
    chrome(s, 4)


def slide_6(prs):
    s = new_slide(prs)
    simple(s, PAD_L, 196, 940, 260,
           "A test that has never failed\nis a test you cannot trust.", 40, INK,
           DISPLAY, line=1.1)
    simple(s, PAD_L, 452, 660, 160,
           "If a brand new test passes straight away, either the work was "
           "already done, or the test is checking nothing. So you always watch "
           "it fail once, on purpose.", 17, INK_SOFT, line=1.45)
    chrome(s, 5)


def slide_7(prs):
    s = new_slide(prs)
    heading(s, "What you get out of it")

    items = [
        ("Every line of code has a test",
         "You only write code to pass a test, so nothing is left untested."),
        ("A safety net that keeps growing",
         "Every old test runs again, so a new change cannot quietly break "
         "old work."),
        ("Bugs are easy to find",
         "A test just broke, and you only changed a few lines. The bug is in "
         "those lines."),
        ("The tests explain the code",
         "Reading the tests shows exactly what the code is supposed to do."),
    ]
    y = 190
    for head, sub in items:
        dash(s, PAD_L, y + 12)
        simple(s, PAD_L + 30, y, 480, 34, head, 17, INK, line=1.25)
        simple(s, PAD_L + 30, y + 28, 480, 60, sub, 13, INK_SOFT, line=1.35)
        y += 108

    stack_figure(s, 872, 226)
    simple(s, 872, 486, 300, 60,
           "Red is the new test. The grey ones are old tests, still running "
           "on every step.", 12, INK_FAINT, line=1.35)
    chrome(s, 6)


def slide_8(prs):
    s = new_slide(prs)
    heading(s, "Where it does not help much")

    cols = [
        [("Screens and visual design",
          "Hard to write an automatic test for \u201cdoes this look right\u201d."),
         ("Big, tangled old systems",
          "You cannot test one small piece on its own.")],
        [("Systems with many things running at once",
          "The same test can pass today and fail tomorrow."),
         ("It only checks what you thought of",
          "So you still need system testing and real users.")],
    ]
    col_w = (CONTENT_W - 56) / 2.0
    for i, col in enumerate(cols):
        x = PAD_L + i * (col_w + 56)
        y = 200
        for head, sub in col:
            dash(s, x, y + 12)
            simple(s, x + 30, y, col_w - 30, 40, head, 17, INK, line=1.25)
            simple(s, x + 30, y + 30, col_w - 30, 70, sub, 13, INK_SOFT,
                   line=1.35)
            y += 130

    chrome(s, 7)


def slide_9(prs):
    s = new_slide(prs)
    heading(s, "Testing at the end vs. testing first")

    cols = [
        ("Testing at the end", INK,
         ["Code first, tests later, if there is time",
          "Bugs show up in a pile, far from where they were made",
          "Hard to say how much of the code is really tested",
          "Changing old code feels risky"]),
        ("Testing first", ACCENT,
         ["One test, one small piece of code, over and over",
          "Bugs show up within minutes of being made",
          "Coverage comes for free",
          "Changing old code is safe, the tests will shout"]),
    ]
    col_w = CONTENT_W / 2.0
    for i, (head, colour, items) in enumerate(cols):
        x = PAD_L + i * col_w
        if i:
            vrule(s, x - 17, 200, 356)
        simple(s, x, 206, col_w - 40, 30, head, 13, colour, bold=True,
               line=1.2)
        y = 252
        for text in items:
            rule(s, x, y + 11, 9, thickness=1, color=INK_FAINT)
            simple(s, x + 22, y, col_w - 52, 60, text, 15, INK_SOFT, line=1.35)
            y += 76

    chrome(s, 8)


def slide_10(prs):
    s = new_slide(prs)
    heading(s, "Remember this")

    items = [
        "Test first, code second, in very small steps.",
        "Red, green, refactor.",
        "It came from agile methods, and it needs an automatic test tool such "
        "as JUnit.",
        "It gives coverage, a safety net, easy debugging and living "
        "documentation.",
        "It replaces none of your system testing or user testing.",
    ]
    y = 196
    for text in items:
        dash(s, PAD_L, y + 12)
        simple(s, PAD_L + 30, y, 900, 60, text, 17, INK, line=1.35)
        y += 84

    chrome(s, 9)


def main():
    prs = Presentation()
    prs.slide_width = px(STAGE_W)
    prs.slide_height = px(STAGE_H)

    core = prs.core_properties
    core.title = "Test-Driven Development \u2014 Unit 8.2, Software Engineering"
    core.subject = "Software Engineering CSC364, Semester VI"
    core.comments = ("Generated from index.html by tools/build_pptx.py. "
                     "Content follows Sommerville, Software Engineering, "
                     "10th edition, chapter 8.")

    for build in (slide_1, slide_2, slide_3, slide_4, slide_5, slide_6,
                  slide_7, slide_8, slide_9, slide_10):
        build(prs)

    prs.save("tdd-presentation.pptx")
    print("wrote tdd-presentation.pptx with %d slides" % len(prs.slides._sldIdLst))


if __name__ == "__main__":
    main()
